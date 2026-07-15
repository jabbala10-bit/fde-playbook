"""
Document extraction + OCR + RAG ingestion pipeline.

Goal:
    Extract content from PDFs, images, and uploaded business documents, chunk the
    text, embed it, store it in a vector index, and query it with provenance.

This file is production-minded but dependency-light:
    - Works out of the box for .txt, .md, .csv, .json, .html.
    - Uses optional packages when available for PDF, OCR, Word, PowerPoint, Excel.
    - Uses deterministic hash embeddings by default so demos/tests run offline.
    - Can switch to OpenAI embeddings with --embedding-provider openai.

Useful optional installs:
    pip install pypdf pdfplumber pillow pytesseract python-docx python-pptx openpyxl beautifulsoup4 openai

OCR note:
    pytesseract also requires the Tesseract executable installed on your OS.

Examples:
    python ai/ocr_code.py --mode guide
    python ai/ocr_code.py --mode ingest --input ./uploads --index ./knowledge_index.json
    python ai/ocr_code.py --mode query --index ./knowledge_index.json --query "What is the refund policy?"
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import html
import json
import logging
import math
import mimetypes
import os
import re
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from textwrap import dedent
from typing import Any, Iterable, Optional, Protocol


LOGGER = logging.getLogger(__name__)

SUPPORTED_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".log",
    ".py",
    ".sql",
    ".json",
    ".jsonl",
    ".csv",
    ".tsv",
    ".html",
    ".htm",
    ".xml",
}
SUPPORTED_DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
SUPPORTED_EXTENSIONS = (
    SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_DOCUMENT_EXTENSIONS | SUPPORTED_IMAGE_EXTENSIONS
)


@dataclass(frozen=True)
class ExtractedDocument:
    """Text extracted from a source file."""

    doc_id: str
    source_path: str
    mime_type: str
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class TextChunk:
    """Chunk ready for embedding."""

    chunk_id: str
    doc_id: str
    source_path: str
    text: str
    chunk_index: int
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class SearchResult:
    """Vector search result with provenance."""

    score: float
    chunk: TextChunk


class EmbeddingProvider(Protocol):
    """Embedding provider contract."""

    dimensions: int

    def embed_texts(self, texts: list[str]) -> list[list[float]]:
        ...


class HashEmbeddingProvider:
    """Offline deterministic embeddings for demos, tests, and local development."""

    def __init__(self, dimensions: int = 384) -> None:
        self.dimensions = dimensions

    def embed_texts(self, texts: list[str]) -> list[list[float]]:
        return [self._embed(text) for text in texts]

    def _embed(self, text: str) -> list[float]:
        vector = [0.0] * self.dimensions
        for token in tokenize(text):
            digest = hashlib.sha256(token.encode("utf-8")).digest()
            index = int.from_bytes(digest[:4], "big") % self.dimensions
            weight = 1.0 + digest[4] / 255.0
            vector[index] += weight

        norm = math.sqrt(sum(value * value for value in vector))
        if norm == 0:
            return vector
        return [value / norm for value in vector]


class OpenAIEmbeddingProvider:
    """OpenAI embedding provider for production-style RAG indexes."""

    def __init__(self, model: str = "text-embedding-3-small") -> None:
        try:
            from openai import OpenAI
        except ImportError as exc:
            raise RuntimeError("Install the OpenAI SDK with: pip install openai") from exc

        self.client = OpenAI()
        self.model = model
        self.dimensions = 1536

    def embed_texts(self, texts: list[str]) -> list[list[float]]:
        response = self.client.embeddings.create(model=self.model, input=texts)
        return [item.embedding for item in response.data]


class DocumentExtractor:
    """Extract text from common uploaded document types."""

    def __init__(self, enable_ocr: bool = True) -> None:
        self.enable_ocr = enable_ocr

    def extract_path(self, path: Path) -> ExtractedDocument:
        path = path.resolve()
        if not path.exists():
            raise FileNotFoundError(path)
        if not path.is_file():
            raise ValueError(f"Expected a file path, got directory: {path}")

        extension = path.suffix.lower()
        mime_type = mimetypes.guess_type(str(path))[0] or "application/octet-stream"

        if extension == ".pdf":
            text, metadata = self._extract_pdf(path)
        elif extension in SUPPORTED_IMAGE_EXTENSIONS:
            text, metadata = self._extract_image(path)
        elif extension == ".docx":
            text, metadata = self._extract_docx(path)
        elif extension == ".pptx":
            text, metadata = self._extract_pptx(path)
        elif extension == ".xlsx":
            text, metadata = self._extract_xlsx(path)
        elif extension in {".csv", ".tsv"}:
            text, metadata = self._extract_delimited(path)
        elif extension in {".json", ".jsonl"}:
            text, metadata = self._extract_json(path)
        elif extension in {".html", ".htm", ".xml"}:
            text, metadata = self._extract_html(path)
        elif extension in SUPPORTED_TEXT_EXTENSIONS:
            text, metadata = self._extract_plain_text(path)
        else:
            raise ValueError(f"Unsupported file extension: {extension}")

        normalized = normalize_text(text)
        return ExtractedDocument(
            doc_id=stable_doc_id(path),
            source_path=str(path),
            mime_type=mime_type,
            text=normalized,
            metadata={
                "filename": path.name,
                "extension": extension,
                "bytes": path.stat().st_size,
                **metadata,
            },
        )

    def _extract_pdf(self, path: Path) -> tuple[str, dict[str, Any]]:
        text_parts: list[str] = []
        metadata: dict[str, Any] = {"extractor": "pdf"}

        try:
            from pypdf import PdfReader
        except ImportError:
            PdfReader = None

        if PdfReader is not None:
            reader = PdfReader(str(path))
            metadata["pages"] = len(reader.pages)
            for page_number, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"[page {page_number}]\n{page_text}")

        if not text_parts:
            text_parts.extend(self._extract_pdf_with_pdfplumber(path, metadata))

        if not text_parts and self.enable_ocr:
            text_parts.extend(self._ocr_pdf_pages(path, metadata))

        return "\n\n".join(text_parts), metadata

    def _extract_pdf_with_pdfplumber(self, path: Path, metadata: dict[str, Any]) -> list[str]:
        try:
            import pdfplumber
        except ImportError:
            return []

        parts: list[str] = []
        metadata["extractor"] = "pdfplumber"
        with pdfplumber.open(str(path)) as pdf:
            metadata["pages"] = len(pdf.pages)
            for page_number, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(f"[page {page_number}]\n{page_text}")
        return parts

    def _ocr_pdf_pages(self, path: Path, metadata: dict[str, Any]) -> list[str]:
        try:
            from pdf2image import convert_from_path
            import pytesseract
        except ImportError:
            metadata["ocr_status"] = "pdf_ocr_dependencies_missing"
            return []

        parts: list[str] = []
        metadata["extractor"] = "pdf_ocr"
        images = convert_from_path(str(path))
        metadata["pages"] = len(images)
        for page_number, image in enumerate(images, start=1):
            page_text = pytesseract.image_to_string(image)
            if page_text.strip():
                parts.append(f"[ocr page {page_number}]\n{page_text}")
        return parts

    def _extract_image(self, path: Path) -> tuple[str, dict[str, Any]]:
        if not self.enable_ocr:
            return "", {"extractor": "image", "ocr_status": "disabled"}

        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            return "", {"extractor": "image", "ocr_status": "dependencies_missing"}

        with Image.open(path) as image:
            text = pytesseract.image_to_string(image)
            metadata = {
                "extractor": "image_ocr",
                "width": image.width,
                "height": image.height,
                "mode": image.mode,
            }
        return text, metadata

    def _extract_docx(self, path: Path) -> tuple[str, dict[str, Any]]:
        try:
            from docx import Document
        except ImportError:
            return self._extract_docx_zip_fallback(path)

        document = Document(str(path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        table_rows: list[str] = []
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    table_rows.append(" | ".join(cells))

        return "\n".join(paragraphs + table_rows), {
            "extractor": "python-docx",
            "paragraphs": len(paragraphs),
            "tables": len(document.tables),
        }

    def _extract_docx_zip_fallback(self, path: Path) -> tuple[str, dict[str, Any]]:
        try:
            with zipfile.ZipFile(path) as archive:
                xml = archive.read("word/document.xml").decode("utf-8", errors="replace")
        except (KeyError, zipfile.BadZipFile):
            return "", {"extractor": "docx_zip_fallback", "status": "failed"}

        text = strip_xml_tags(xml)
        return text, {"extractor": "docx_zip_fallback"}

    def _extract_pptx(self, path: Path) -> tuple[str, dict[str, Any]]:
        try:
            from pptx import Presentation
        except ImportError:
            return "", {"extractor": "pptx", "status": "python-pptx_missing"}

        presentation = Presentation(str(path))
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                parts.append(f"[slide {slide_index}]\n" + "\n".join(slide_text))

        return "\n\n".join(parts), {
            "extractor": "python-pptx",
            "slides": len(presentation.slides),
        }

    def _extract_xlsx(self, path: Path) -> tuple[str, dict[str, Any]]:
        try:
            from openpyxl import load_workbook
        except ImportError:
            return "", {"extractor": "xlsx", "status": "openpyxl_missing"}

        workbook = load_workbook(str(path), data_only=True, read_only=True)
        parts: list[str] = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [format_cell(value) for value in row]
                if any(values):
                    rows.append(" | ".join(values))
            if rows:
                parts.append(f"[sheet {sheet_name}]\n" + "\n".join(rows))

        return "\n\n".join(parts), {
            "extractor": "openpyxl",
            "sheets": workbook.sheetnames,
        }

    def _extract_delimited(self, path: Path) -> tuple[str, dict[str, Any]]:
        delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
        rows: list[str] = []
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            reader = csv.reader(handle, delimiter=delimiter)
            for row in reader:
                rows.append(" | ".join(cell.strip() for cell in row))
        return "\n".join(rows), {"extractor": "csv", "rows": len(rows)}

    def _extract_json(self, path: Path) -> tuple[str, dict[str, Any]]:
        raw = path.read_text(encoding="utf-8-sig")
        if path.suffix.lower() == ".jsonl":
            records = [json.loads(line) for line in raw.splitlines() if line.strip()]
            return "\n".join(json.dumps(record, ensure_ascii=False) for record in records), {
                "extractor": "jsonl",
                "records": len(records),
            }

        payload = json.loads(raw)
        return json.dumps(payload, ensure_ascii=False, indent=2), {"extractor": "json"}

    def _extract_html(self, path: Path) -> tuple[str, dict[str, Any]]:
        raw = path.read_text(encoding="utf-8-sig", errors="replace")
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            return strip_html_tags(raw), {"extractor": "html_regex"}

        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return soup.get_text("\n"), {"extractor": "beautifulsoup"}

    def _extract_plain_text(self, path: Path) -> tuple[str, dict[str, Any]]:
        return path.read_text(encoding="utf-8-sig", errors="replace"), {"extractor": "plain_text"}


class InMemoryVectorStore:
    """Simple vector store with JSON persistence."""

    def __init__(self, dimensions: int) -> None:
        self.dimensions = dimensions
        self.rows: list[tuple[TextChunk, list[float]]] = []

    def add(self, chunks: list[TextChunk], embeddings: list[list[float]]) -> None:
        if len(chunks) != len(embeddings):
            raise ValueError("chunks and embeddings must have the same length")

        for chunk, embedding in zip(chunks, embeddings):
            if len(embedding) != self.dimensions:
                raise ValueError(
                    f"Embedding dimension mismatch: expected {self.dimensions}, got {len(embedding)}"
                )
            self.rows.append((chunk, embedding))

    def search(
        self,
        query: str,
        embedding_provider: EmbeddingProvider,
        *,
        top_k: int = 5,
        filters: Optional[dict[str, Any]] = None,
    ) -> list[SearchResult]:
        query_vector = embedding_provider.embed_texts([query])[0]
        results: list[SearchResult] = []

        for chunk, vector in self.rows:
            if filters and not metadata_matches(chunk.metadata, filters):
                continue
            score = cosine_similarity(query_vector, vector)
            results.append(SearchResult(score=score, chunk=chunk))

        results.sort(key=lambda result: result.score, reverse=True)
        return results[:top_k]

    def save(self, path: Path) -> None:
        payload = {
            "dimensions": self.dimensions,
            "rows": [
                {
                    "chunk": asdict(chunk),
                    "embedding": embedding,
                }
                for chunk, embedding in self.rows
            ],
        }
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(payload, indent=2), encoding="utf-8")

    @classmethod
    def load(cls, path: Path) -> "InMemoryVectorStore":
        payload = json.loads(path.read_text(encoding="utf-8"))
        store = cls(dimensions=int(payload["dimensions"]))
        for row in payload.get("rows", []):
            chunk_payload = row["chunk"]
            chunk = TextChunk(
                chunk_id=chunk_payload["chunk_id"],
                doc_id=chunk_payload["doc_id"],
                source_path=chunk_payload["source_path"],
                text=chunk_payload["text"],
                chunk_index=int(chunk_payload["chunk_index"]),
                metadata=dict(chunk_payload.get("metadata", {})),
            )
            store.rows.append((chunk, [float(value) for value in row["embedding"]]))
        return store


class DocumentRAGPipeline:
    """End-to-end extraction, chunking, embedding, storage, and retrieval."""

    def __init__(
        self,
        extractor: Optional[DocumentExtractor] = None,
        embedding_provider: Optional[EmbeddingProvider] = None,
        chunk_size: int = 900,
        chunk_overlap: int = 150,
    ) -> None:
        if chunk_size <= 0:
            raise ValueError("chunk_size must be positive")
        if chunk_overlap < 0 or chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap must be >= 0 and < chunk_size")

        self.extractor = extractor or DocumentExtractor()
        self.embedding_provider = embedding_provider or HashEmbeddingProvider()
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.store = InMemoryVectorStore(self.embedding_provider.dimensions)

    def ingest_paths(self, paths: Iterable[Path]) -> list[ExtractedDocument]:
        documents: list[ExtractedDocument] = []
        chunks: list[TextChunk] = []

        for path in expand_input_paths(paths):
            try:
                document = self.extractor.extract_path(path)
            except Exception as exc:
                LOGGER.warning("Skipping %s: %s", path, exc)
                continue

            if not document.text.strip():
                LOGGER.warning("No extractable text found in %s", path)
                continue

            documents.append(document)
            chunks.extend(
                chunk_document(
                    document,
                    chunk_size=self.chunk_size,
                    overlap=self.chunk_overlap,
                )
            )

        if chunks:
            embeddings = self.embedding_provider.embed_texts([chunk.text for chunk in chunks])
            self.store.add(chunks, embeddings)

        return documents

    def query(self, question: str, top_k: int = 5) -> dict[str, Any]:
        results = self.store.search(question, self.embedding_provider, top_k=top_k)
        context = "\n\n".join(
            f"[{result.chunk.chunk_id}] {result.chunk.text}" for result in results
        )
        answer = rule_based_answer(question, results)
        return {
            "question": question,
            "answer": answer,
            "context": context,
            "sources": [
                {
                    "chunk_id": result.chunk.chunk_id,
                    "doc_id": result.chunk.doc_id,
                    "source_path": result.chunk.source_path,
                    "score": round(result.score, 4),
                    "metadata": result.chunk.metadata,
                    "preview": result.chunk.text[:300],
                }
                for result in results
            ],
        }


def chunk_document(document: ExtractedDocument, chunk_size: int, overlap: int) -> list[TextChunk]:
    words = document.text.split()
    if not words:
        return []

    chunks: list[TextChunk] = []
    step = chunk_size - overlap
    for index, start in enumerate(range(0, len(words), step)):
        selected = words[start : start + chunk_size]
        if not selected:
            break
        text = " ".join(selected)
        chunk_id = f"{document.doc_id}::chunk::{index}"
        chunks.append(
            TextChunk(
                chunk_id=chunk_id,
                doc_id=document.doc_id,
                source_path=document.source_path,
                text=text,
                chunk_index=index,
                metadata={
                    **document.metadata,
                    "word_start": start,
                    "word_end": start + len(selected),
                },
            )
        )
    return chunks


def expand_input_paths(paths: Iterable[Path]) -> list[Path]:
    expanded: list[Path] = []
    for path in paths:
        if path.is_dir():
            for child in sorted(path.rglob("*")):
                if child.is_file() and child.suffix.lower() in SUPPORTED_EXTENSIONS:
                    expanded.append(child)
        elif path.is_file():
            expanded.append(path)
        else:
            LOGGER.warning("Input path does not exist: %s", path)
    return expanded


def build_embedding_provider(provider: str, model: Optional[str] = None) -> EmbeddingProvider:
    provider = provider.lower().strip()
    if provider == "hash":
        return HashEmbeddingProvider()
    if provider == "openai":
        return OpenAIEmbeddingProvider(model=model or "text-embedding-3-small")
    raise ValueError("embedding provider must be 'hash' or 'openai'")


def rule_based_answer(question: str, results: list[SearchResult]) -> str:
    if not results:
        return "No relevant context was found in the indexed documents."

    citations = ", ".join(result.chunk.chunk_id for result in results[:3])
    return (
        "I found relevant context in the knowledge index. Use the retrieved chunks "
        f"to answer the question: {question!r}. Top citations: {citations}."
    )


def tokenize(text: str) -> list[str]:
    return re.findall(r"[a-zA-Z0-9_]+", text.lower())


def normalize_text(text: str) -> str:
    text = html.unescape(text)
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def stable_doc_id(path: Path) -> str:
    resolved = str(path.resolve()).lower()
    digest = hashlib.sha1(resolved.encode("utf-8")).hexdigest()[:12]
    return f"{path.stem}-{digest}"


def cosine_similarity(left: list[float], right: list[float]) -> float:
    if len(left) != len(right):
        raise ValueError("Vectors must have the same dimension")
    dot = sum(a * b for a, b in zip(left, right))
    left_norm = math.sqrt(sum(a * a for a in left))
    right_norm = math.sqrt(sum(b * b for b in right))
    if left_norm == 0 or right_norm == 0:
        return 0.0
    return dot / (left_norm * right_norm)


def metadata_matches(metadata: dict[str, Any], filters: dict[str, Any]) -> bool:
    for key, expected in filters.items():
        if metadata.get(key) != expected:
            return False
    return True


def strip_html_tags(raw: str) -> str:
    return normalize_text(re.sub(r"<[^>]+>", " ", raw))


def strip_xml_tags(raw: str) -> str:
    raw = re.sub(r"<w:tab\s*/>", "\t", raw)
    raw = re.sub(r"</w:p>", "\n", raw)
    return strip_html_tags(raw)


def format_cell(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def describe_architecture() -> str:
    return dedent(
        """
        Document-to-RAG architecture:

        1. Upload/Input
           PDFs, scanned PDFs, images, DOCX, PPTX, XLSX, CSV, JSON, HTML, TXT.

        2. Extraction
           Native text extraction first. OCR fallback for images and scanned PDFs.

        3. Normalization
           Clean whitespace, preserve page/sheet/slide markers, attach metadata.

        4. Chunking
           Split by word windows with overlap. Keep source path, document ID,
           chunk index, and positional metadata.

        5. Embedding
           Use hash embeddings for offline demos or OpenAI embeddings for semantic retrieval.

        6. Vector Store
           Store chunk + embedding + provenance. This demo persists JSON; production
           systems can swap in pgvector, Pinecone, Weaviate, Milvus, or FAISS.

        7. Retrieval/RAG
           Embed the query, retrieve top chunks, build context, answer with citations.

        Production upgrades:
           Add file deduplication, async OCR queues, document-level ACLs, PII redaction,
           embedding batch retries, index versioning, hybrid lexical/vector search,
           reranking, and citation verification.
        """
    ).strip()


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract uploaded docs and build a RAG knowledge index")
    parser.add_argument(
        "--mode",
        choices=("guide", "ingest", "query"),
        default="guide",
        help="Run guide, ingest documents, or query an existing index.",
    )
    parser.add_argument(
        "--input",
        action="append",
        default=[],
        help="File or directory to ingest. Can be provided multiple times.",
    )
    parser.add_argument("--index", default="knowledge_index.json", help="Path to index JSON file.")
    parser.add_argument("--query", default="", help="Question for query mode.")
    parser.add_argument("--top-k", type=int, default=5)
    parser.add_argument("--chunk-size", type=int, default=900)
    parser.add_argument("--chunk-overlap", type=int, default=150)
    parser.add_argument("--disable-ocr", action="store_true")
    parser.add_argument("--embedding-provider", choices=("hash", "openai"), default="hash")
    parser.add_argument("--embedding-model", default=None)
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(name)s %(message)s")

    if args.mode == "guide":
        print(describe_architecture())
        return

    embedding_provider = build_embedding_provider(args.embedding_provider, args.embedding_model)
    index_path = Path(args.index)

    if args.mode == "ingest":
        if not args.input:
            raise ValueError("--input is required for ingest mode")

        pipeline = DocumentRAGPipeline(
            extractor=DocumentExtractor(enable_ocr=not args.disable_ocr),
            embedding_provider=embedding_provider,
            chunk_size=args.chunk_size,
            chunk_overlap=args.chunk_overlap,
        )
        documents = pipeline.ingest_paths([Path(value) for value in args.input])
        pipeline.store.save(index_path)
        print(
            json.dumps(
                {
                    "documents_ingested": len(documents),
                    "chunks_indexed": len(pipeline.store.rows),
                    "index_path": str(index_path.resolve()),
                },
                indent=2,
            )
        )
        return

    if not args.query:
        raise ValueError("--query is required for query mode")

    store = InMemoryVectorStore.load(index_path)
    if store.dimensions != embedding_provider.dimensions:
        raise ValueError(
            "Index embedding dimensions do not match selected embedding provider. "
            "Use the same provider/model used during ingestion."
        )
    pipeline = DocumentRAGPipeline(embedding_provider=embedding_provider)
    pipeline.store = store
    print(json.dumps(pipeline.query(args.query, top_k=args.top_k), indent=2))


if __name__ == "__main__":
    main()
